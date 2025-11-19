"""
Document Input Adapter - Handles PDFs, DOCX, presentations, emails
"""
from __future__ import annotations

import asyncio
import logging
import time
from typing import Any, Dict, List, Optional
import io
import zipfile
import re

from .base import BaseInputAdapter, InputType, ProcessedInput, InputMetadata, ProcessingResult

# Optional imports with fallbacks
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import email
    from email.parser import Parser as EmailParser
except ImportError:
    email = None
    EmailParser = None

log = logging.getLogger("document-adapter")

class DocumentAdapter(BaseInputAdapter):
    """Adapter for document processing (PDF, DOCX, PPTX, emails, HTML)"""
    
    def __init__(self):
        super().__init__("DocumentAdapter")
        
    async def can_handle(self, input_data: Any, input_type: InputType, metadata: InputMetadata) -> bool:
        """Check if can handle document input"""
        if input_type in [InputType.PDF, InputType.DOCX, InputType.PPTX, InputType.HTML, InputType.RTF]:
            return True
            
        # Check by filename
        if metadata.filename:
            filename = metadata.filename.lower()
            if filename.endswith(('.pdf', '.docx', '.pptx', '.html', '.htm', '.rtf')):
                return True
                
        # Check by content type
        if metadata.content_type:
            content_type = metadata.content_type.lower()
            if any(doc_type in content_type for doc_type in [
                'pdf', 'word', 'powerpoint', 'html', 'rtf'
            ]):
                return True
                
        # Check binary signatures
        if isinstance(input_data, bytes):
            if input_data.startswith(b'%PDF'):  # PDF
                return True
            elif input_data.startswith(b'PK\x03\x04'):  # ZIP-based (DOCX, PPTX)
                return True
            elif input_data.startswith(b'<html') or input_data.startswith(b'<!DOCTYPE'):  # HTML
                return True
                
        return False
        
    async def process(self, input_data: Any, metadata: InputMetadata) -> ProcessedInput:
        """Process document input"""
        start_time = time.time()
        
        try:
            # Detect document type
            doc_type = self._detect_document_type(input_data, metadata)
            
            # Process based on type
            if doc_type == InputType.PDF:
                result = await self._process_pdf(input_data, metadata)
            elif doc_type == InputType.DOCX:
                result = await self._process_docx(input_data, metadata)
            elif doc_type == InputType.PPTX:
                result = await self._process_pptx(input_data, metadata)
            elif doc_type == InputType.HTML:
                result = await self._process_html(input_data, metadata)
            else:
                result = await self._process_generic_document(input_data, metadata, doc_type)
                
            processing_time = time.time() - start_time
            result.processing_time = processing_time
            
            self.update_stats(processing_time, result.result_status == ProcessingResult.SUCCESS)
            return result
            
        except Exception as e:
            processing_time = time.time() - start_time
            log.error(f"Document processing failed: {e}")
            
            result = ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=InputType.PDF,  # Default
                processed_type="document",
                content="",
                metadata=metadata.to_dict(),
                processing_time=processing_time,
                result_status=ProcessingResult.FAILED,
                error_message=str(e)
            )
            
            self.update_stats(processing_time, False)
            return result
            
    def get_supported_formats(self) -> List[InputType]:
        """Get supported document formats"""
        return [InputType.PDF, InputType.DOCX, InputType.PPTX, InputType.HTML, InputType.RTF]
        
    def _detect_document_type(self, input_data: Any, metadata: InputMetadata) -> InputType:
        """Detect specific document type"""
        # Check filename first
        if metadata.filename:
            filename = metadata.filename.lower()
            if filename.endswith('.pdf'):
                return InputType.PDF
            elif filename.endswith('.docx'):
                return InputType.DOCX
            elif filename.endswith('.pptx'):
                return InputType.PPTX
            elif filename.endswith(('.html', '.htm')):
                return InputType.HTML
            elif filename.endswith('.rtf'):
                return InputType.RTF
                
        # Check content type
        if metadata.content_type:
            content_type = metadata.content_type.lower()
            if 'pdf' in content_type:
                return InputType.PDF
            elif 'word' in content_type or 'officedocument.wordprocessing' in content_type:
                return InputType.DOCX
            elif 'powerpoint' in content_type or 'officedocument.presentation' in content_type:
                return InputType.PPTX
            elif 'html' in content_type:
                return InputType.HTML
                
        # Check binary signatures
        if isinstance(input_data, bytes):
            if input_data.startswith(b'%PDF'):
                return InputType.PDF
            elif input_data.startswith(b'PK\x03\x04'):
                # Could be DOCX or PPTX - check internal structure
                try:
                    with zipfile.ZipFile(io.BytesIO(input_data), 'r') as zip_file:
                        file_list = zip_file.namelist()
                        if 'word/document.xml' in file_list:
                            return InputType.DOCX
                        elif 'ppt/presentation.xml' in file_list:
                            return InputType.PPTX
                except:
                    pass
            elif input_data.startswith((b'<html', b'<!DOCTYPE')):
                return InputType.HTML
                
        # Check if it's text that looks like HTML
        if isinstance(input_data, str):
            if input_data.strip().startswith(('<html', '<!DOCTYPE')):
                return InputType.HTML
                
        return InputType.PDF  # Default fallback
        
    async def _process_pdf(self, input_data: Any, metadata: InputMetadata) -> ProcessedInput:
        """Process PDF document"""
        if not PdfReader:
            return await self._process_generic_document(input_data, metadata, InputType.PDF)
            
        try:
            # Ensure we have bytes
            if isinstance(input_data, str):
                pdf_bytes = input_data.encode('latin-1')
            else:
                pdf_bytes = input_data
                
            # Read PDF
            pdf_file = io.BytesIO(pdf_bytes)
            pdf_reader = PdfReader(pdf_file)
            
            # Extract text from all pages
            text_content = ""
            page_count = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    text_content += f"\\n--- Page {page_num + 1} ---\\n{page_text}\\n"
                except Exception as e:
                    log.debug(f"Could not extract text from page {page_num + 1}: {e}")
                    
            # Extract metadata
            pdf_metadata = {}
            if hasattr(pdf_reader, 'metadata') and pdf_reader.metadata:
                for key, value in pdf_reader.metadata.items():
                    if value:
                        pdf_metadata[key.replace('/', '')] = str(value)
                        
            # Extract features
            features = self._extract_basic_features(text_content, metadata)
            features.update({
                "document_type": "PDF",
                "page_count": page_count,
                "has_metadata": bool(pdf_metadata),
                "metadata_fields": list(pdf_metadata.keys()),
                "estimated_word_count": len(text_content.split()),
                "has_images": self._pdf_has_images(pdf_reader),
                "is_searchable": len(text_content.strip()) > 0
            })
            
            # Combine content
            content = {
                "text": text_content,
                "metadata": pdf_metadata,
                "pages": page_count
            }
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=InputType.PDF,
                processed_type="document",
                content=content,
                metadata=metadata.to_dict(),
                extracted_features=features,
                confidence=0.9 if text_content.strip() else 0.5
            )
            
        except Exception as e:
            log.error(f"PDF processing failed: {e}")
            return await self._process_generic_document(input_data, metadata, InputType.PDF)
            
    async def _process_docx(self, input_data: Any, metadata: InputMetadata) -> ProcessedInput:
        """Process DOCX document"""
        if not DocxDocument:
            return await self._process_generic_document(input_data, metadata, InputType.DOCX)
            
        try:
            # Ensure we have bytes
            if isinstance(input_data, str):
                docx_bytes = input_data.encode('latin-1')
            else:
                docx_bytes = input_data
                
            # Read DOCX
            docx_file = io.BytesIO(docx_bytes)
            document = DocxDocument(docx_file)
            
            # Extract text from paragraphs
            text_content = ""
            paragraph_count = 0
            
            for paragraph in document.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\\n"
                    paragraph_count += 1
                    
            # Extract text from tables
            table_count = 0
            for table in document.tables:
                table_count += 1
                text_content += f"\\n--- Table {table_count} ---\\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text_content += row_text + "\\n"
                    
            # Extract document properties
            doc_properties = {}
            if hasattr(document, 'core_properties'):
                props = document.core_properties
                for attr in ['author', 'title', 'subject', 'keywords', 'created', 'modified']:
                    value = getattr(props, attr, None)
                    if value:
                        doc_properties[attr] = str(value)
                        
            # Extract features
            features = self._extract_basic_features(text_content, metadata)
            features.update({
                "document_type": "DOCX",
                "paragraph_count": paragraph_count,
                "table_count": table_count,
                "has_properties": bool(doc_properties),
                "property_fields": list(doc_properties.keys()),
                "estimated_word_count": len(text_content.split())
            })
            
            # Combine content
            content = {
                "text": text_content,
                "properties": doc_properties,
                "structure": {
                    "paragraphs": paragraph_count,
                    "tables": table_count
                }
            }
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=InputType.DOCX,
                processed_type="document",
                content=content,
                metadata=metadata.to_dict(),
                extracted_features=features,
                confidence=0.9
            )
            
        except Exception as e:
            log.error(f"DOCX processing failed: {e}")
            return await self._process_generic_document(input_data, metadata, InputType.DOCX)
            
    async def _process_pptx(self, input_data: Any, metadata: InputMetadata) -> ProcessedInput:
        """Process PPTX presentation"""
        if not Presentation:
            return await self._process_generic_document(input_data, metadata, InputType.PPTX)
            
        try:
            # Ensure we have bytes
            if isinstance(input_data, str):
                pptx_bytes = input_data.encode('latin-1')
            else:
                pptx_bytes = input_data
                
            # Read PPTX
            pptx_file = io.BytesIO(pptx_bytes)
            presentation = Presentation(pptx_file)
            
            # Extract text from slides
            text_content = ""
            slide_count = len(presentation.slides)
            
            for slide_num, slide in enumerate(presentation.slides):
                text_content += f"\\n--- Slide {slide_num + 1} ---\\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\\n"
                        
            # Extract features
            features = self._extract_basic_features(text_content, metadata)
            features.update({
                "document_type": "PPTX",
                "slide_count": slide_count,
                "estimated_word_count": len(text_content.split()),
                "has_content": len(text_content.strip()) > 0
            })
            
            # Combine content
            content = {
                "text": text_content,
                "slides": slide_count
            }
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=InputType.PPTX,
                processed_type="document",
                content=content,
                metadata=metadata.to_dict(),
                extracted_features=features,
                confidence=0.9
            )
            
        except Exception as e:
            log.error(f"PPTX processing failed: {e}")
            return await self._process_generic_document(input_data, metadata, InputType.PPTX)
            
    async def _process_html(self, input_data: Any, metadata: InputMetadata) -> ProcessedInput:
        """Process HTML document"""
        try:
            # Convert to string if needed
            if isinstance(input_data, bytes):
                html_content = input_data.decode('utf-8', errors='replace')
            else:
                html_content = str(input_data)
                
            # Extract text content
            if BeautifulSoup:
                # Use BeautifulSoup for better HTML parsing
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                    
                # Extract text
                text_content = soup.get_text()
                
                # Extract metadata
                html_metadata = {}
                
                # Title
                title = soup.find('title')
                if title:
                    html_metadata['title'] = title.get_text().strip()
                    
                # Meta tags
                meta_tags = soup.find_all('meta')
                for meta in meta_tags:
                    name = meta.get('name') or meta.get('property')
                    content = meta.get('content')
                    if name and content:
                        html_metadata[name] = content
                        
                # Links
                links = soup.find_all('a', href=True)
                html_metadata['link_count'] = len(links)
                
                # Images
                images = soup.find_all('img', src=True)
                html_metadata['image_count'] = len(images)
                
            else:
                # Fallback: simple regex-based extraction
                text_content = self._extract_text_from_html(html_content)
                html_metadata = self._extract_html_metadata(html_content)
                
            # Clean up text
            text_content = re.sub(r'\\s+', ' ', text_content).strip()
            
            # Extract features
            features = self._extract_basic_features(text_content, metadata)
            features.update({
                "document_type": "HTML",
                "has_title": 'title' in html_metadata,
                "has_meta_tags": len([k for k in html_metadata.keys() if k not in ['link_count', 'image_count']]) > 0,
                "link_count": html_metadata.get('link_count', 0),
                "image_count": html_metadata.get('image_count', 0),
                "estimated_word_count": len(text_content.split())
            })
            
            # Combine content
            content = {
                "text": text_content,
                "html_metadata": html_metadata,
                "raw_html": html_content[:1000] + "..." if len(html_content) > 1000 else html_content
            }
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=InputType.HTML,
                processed_type="document",
                content=content,
                metadata=metadata.to_dict(),
                extracted_features=features,
                confidence=0.8
            )
            
        except Exception as e:
            log.error(f"HTML processing failed: {e}")
            return await self._process_generic_document(input_data, metadata, InputType.HTML)
            
    async def _process_generic_document(self, input_data: Any, metadata: InputMetadata, doc_type: InputType) -> ProcessedInput:
        """Generic document processing fallback"""
        try:
            # Convert to text
            if isinstance(input_data, bytes):
                text_content = input_data.decode('utf-8', errors='replace')
            else:
                text_content = str(input_data)
                
            # Extract basic features
            features = self._extract_basic_features(text_content, metadata)
            features.update({
                "document_type": doc_type.value,
                "processing_method": "generic",
                "estimated_word_count": len(text_content.split())
            })
            
            content = {
                "text": text_content,
                "processing_note": f"Processed as generic {doc_type.value} - specialized parser not available"
            }
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=doc_type,
                processed_type="document",
                content=content,
                metadata=metadata.to_dict(),
                extracted_features=features,
                result_status=ProcessingResult.PARTIAL,
                confidence=0.6
            )
            
        except Exception as e:
            log.error(f"Generic document processing failed: {e}")
            
            return ProcessedInput(
                input_id=self._generate_input_id(input_data, metadata),
                original_type=doc_type,
                processed_type="document",
                content="",
                metadata=metadata.to_dict(),
                result_status=ProcessingResult.FAILED,
                error_message=str(e),
                confidence=0.0
            )
            
    def _pdf_has_images(self, pdf_reader) -> bool:
        """Check if PDF contains images"""
        try:
            for page in pdf_reader.pages:
                if hasattr(page, 'images') and len(page.images) > 0:
                    return True
                # Alternative check for images in page resources
                if '/XObject' in page.get('/Resources', {}):
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        if xobjects[obj]['/Subtype'] == '/Image':
                            return True
        except Exception:
            pass
        return False
        
    def _extract_text_from_html(self, html_content: str) -> str:
        """Extract text from HTML using regex (fallback method)"""
        # Remove script and style content
        html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        
        # Remove HTML tags
        text_content = re.sub(r'<[^>]+>', ' ', html_content)
        
        # Decode HTML entities
        import html
        text_content = html.unescape(text_content)
        
        return text_content
        
    def _extract_html_metadata(self, html_content: str) -> Dict[str, Any]:
        """Extract metadata from HTML using regex (fallback method)"""
        metadata = {}
        
        # Extract title
        title_match = re.search(r'<title[^>]*>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
        if title_match:
            metadata['title'] = title_match.group(1).strip()
            
        # Count links and images
        link_count = len(re.findall(r'<a[^>]+href=', html_content, re.IGNORECASE))
        image_count = len(re.findall(r'<img[^>]+src=', html_content, re.IGNORECASE))
        
        metadata['link_count'] = link_count
        metadata['image_count'] = image_count
        
        return metadata
